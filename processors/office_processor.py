from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import os
from processors.base_processor import BaseProcessor

class OfficeProcessor(BaseProcessor):
    """Xử lý file Office (DOCX, XLSX, PPTX)"""
    
    def extract_text(self) -> str:
        """Trích xuất văn bản từ file Office"""
        ext = os.path.splitext(self.file_path)[1].lower()
        
        if ext == '.docx':
            return self._extract_docx()
        elif ext == '.xlsx':
            return self._extract_xlsx()
        elif ext == '.pptx':
            return self._extract_pptx()
        else:
            raise ValueError(f"Định dạng Office không được hỗ trợ: {ext}")
    
    def _extract_docx(self) -> str:
        """Trích xuất từ DOCX"""
        text_parts = []
        
        try:
            doc = Document(self.file_path)
            
            # Lấy văn bản từ paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Lấy văn bản từ tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_parts.append("\t".join(row_text))
            
        except Exception as e:
            return f"Lỗi xử lý DOCX: {str(e)}"
        
        return "\n".join(text_parts)
    
    def _extract_xlsx(self) -> str:
        """Trích xuất từ XLSX"""
        text_parts = []
        
        try:
            wb = load_workbook(self.file_path, data_only=True)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            row_data.append(str(cell))
                    if any(row_data):
                        text_parts.append("\t".join(row_data))
                
        except Exception as e:
            return f"Lỗi xử lý XLSX: {str(e)}"
        
        return "\n".join(text_parts)
    
    def _extract_pptx(self) -> str:
        """Trích xuất từ PPTX"""
        text_parts = []
        
        try:
            prs = Presentation(self.file_path)
            
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"--- Slide {i+1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
        except Exception as e:
            return f"Lỗi xử lý PPTX: {str(e)}"
        
        return "\n".join(text_parts)